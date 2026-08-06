#!/usr/bin/env python3
"""Generate Argos PPT for AMD Radeon Hackathon Track 2 submission."""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# Colors
AMD_RED = RGBColor(0xED, 0x1C, 0x24)
DARK_BG = RGBColor(0x1A, 0x1A, 0x2E)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY = RGBColor(0xCC, 0xCC, 0xCC)
CYAN = RGBColor(0x00, 0xCF, 0xFF)
GOLD = RGBColor(0xFF, 0xD7, 0x00)
GREEN = RGBColor(0x00, 0xE6, 0x76)


def add_bg(slide, color=DARK_BG):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_textbox(slide, left, top, width, height, text, font_size=18,
                color=WHITE, bold=False, alignment=PP_ALIGN.LEFT, font_name="Segoe UI"):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    return txBox


def add_bullet_list(slide, left, top, width, height, items, font_size=16,
                    color=LIGHT_GRAY, font_name="Segoe UI"):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = item
        p.font.size = Pt(font_size)
        p.font.color.rgb = color
        p.font.name = font_name
        p.space_after = Pt(6)
    return txBox


def add_shape_with_text(slide, left, top, width, height, text,
                        fill_color, text_color=WHITE, font_size=14, bold=True):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.color.rgb = fill_color
    tf = shape.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = text_color
    p.font.bold = bold
    p.font.name = "Segoe UI"
    p.alignment = PP_ALIGN.CENTER
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    return shape


# ── Slide 1: Title ──
slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank
add_bg(slide)
add_textbox(slide, Inches(1), Inches(0.8), Inches(11), Inches(0.6),
            "AMD Radeon Hackathon 2026", 20, AMD_RED, True, PP_ALIGN.CENTER)
add_textbox(slide, Inches(1), Inches(1.8), Inches(11), Inches(1.2),
            "Argos", 54, WHITE, True, PP_ALIGN.CENTER)
add_textbox(slide, Inches(1), Inches(3.0), Inches(11), Inches(0.8),
            "Agentic AI Companion with Local RAG, Tool Invocation & Voice", 24, CYAN, False, PP_ALIGN.CENTER)
add_textbox(slide, Inches(1), Inches(4.2), Inches(11), Inches(0.5),
            "Track 2: Development & Local Deployment of Private AI Agents", 18, LIGHT_GRAY, False, PP_ALIGN.CENTER)
add_textbox(slide, Inches(1), Inches(5.0), Inches(11), Inches(0.5),
            "Team: GERALD BUSTILLA", 20, GOLD, True, PP_ALIGN.CENTER)
add_textbox(slide, Inches(1), Inches(5.8), Inches(11), Inches(0.5),
            "Powered by AMD Radeon Cloud API (OpenAI-compatible) + Local Whisper STT", 16, GREEN, False, PP_ALIGN.CENTER)

# ── Slide 2: Application Scenarios ──
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_textbox(slide, Inches(0.5), Inches(0.3), Inches(12), Inches(0.7),
            "Application Scenarios", 32, AMD_RED, True)

scenarios = [
    ("Personal Intelligent Assistant",
     "Always-on desktop companion with real-time chat, proactive engagement,\nvoice input (Whisper), and screen awareness"),
    ("Office Automation Assistant",
     "Browser automation, file management, command execution,\nkeyboard/mouse control via 38+ tools"),
    ("Local Knowledge Base Q&A (RAG)",
     "TF-IDF + BM25 document indexing with Reciprocal Rank Fusion.\nManual folder sync, persistent index, 100% local search"),
    ("Privacy-First Design",
     "No hardcoded API keys, configurable endpoints, no telemetry.\nAll RAG, memory, and tool execution runs locally"),
]

y = Inches(1.3)
for title, desc in scenarios:
    add_shape_with_text(slide, Inches(0.5), y, Inches(3.5), Inches(1.2),
                        title, AMD_RED, WHITE, 14)
    add_textbox(slide, Inches(4.3), y + Inches(0.15), Inches(8.5), Inches(1.0),
                desc, 16, LIGHT_GRAY)
    y += Inches(1.4)

# ── Slide 3: Architecture ──
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_textbox(slide, Inches(0.5), Inches(0.3), Inches(12), Inches(0.7),
            "Agent Architecture", 32, AMD_RED, True)

# Architecture boxes
add_shape_with_text(slide, Inches(0.5), Inches(1.3), Inches(3.5), Inches(0.8),
                    "Robot Overlay\n(Direct2D)", RGBColor(0x2A, 0x2A, 0x4A), CYAN, 12)
add_shape_with_text(slide, Inches(4.3), Inches(1.3), Inches(3.5), Inches(0.8),
                    "Speech Bubble\n(Chat UI + Settings)", RGBColor(0x2A, 0x2A, 0x4A), CYAN, 12)
add_shape_with_text(slide, Inches(8.1), Inches(1.3), Inches(4.5), Inches(0.8),
                    "System Tray\n(Hide/Show/Quit)", RGBColor(0x2A, 0x2A, 0x4A), CYAN, 12)

add_shape_with_text(slide, Inches(1.5), Inches(2.6), Inches(10), Inches(0.6),
                    "Agent Core (C++17)", RGBColor(0x0D, 0x47, 0x47), WHITE, 16)

add_shape_with_text(slide, Inches(0.5), Inches(3.5), Inches(2.8), Inches(1.0),
                    "AgentClient\nSSE Streaming\nTool Loop (8 iter)", RGBColor(0x1B, 0x5E, 0x20), GREEN, 11)
add_shape_with_text(slide, Inches(3.6), Inches(3.5), Inches(2.8), Inches(1.0),
                    "Tool Dispatcher\n38+ Tools\nFile/Browser/Screen", RGBColor(0x1B, 0x5E, 0x20), GREEN, 11)
add_shape_with_text(slide, Inches(6.7), Inches(3.5), Inches(2.8), Inches(1.0),
                    "RAG Engine\nTF-IDF + BM25\nRRF Fusion", RGBColor(0x1B, 0x5E, 0x20), GREEN, 11)
add_shape_with_text(slide, Inches(9.8), Inches(3.5), Inches(3.0), Inches(1.0),
                    "Memory (JSONL)\n+ Whisper STT\n(Local, GGML)", RGBColor(0x1B, 0x5E, 0x20), GREEN, 11)

add_shape_with_text(slide, Inches(2.5), Inches(5.0), Inches(8), Inches(0.8),
                    "AMD Radeon Cloud API (OpenAI-compatible)\nDeepSeek-V4-Flash | MiniCPM5-1B | Qwen3.6-35B-A3B",
                    AMD_RED, WHITE, 13)

add_textbox(slide, Inches(0.5), Inches(6.2), Inches(12), Inches(0.8),
            "Local: Whisper STT, RAG, Memory, All Tools    |    Cloud: LLM Inference on AMD Radeon GPU",
            14, GOLD, True, PP_ALIGN.CENTER)

# ── Slide 4: Core Capabilities ──
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_textbox(slide, Inches(0.5), Inches(0.3), Inches(12), Inches(0.7),
            "Core Capabilities (5 of 5 Implemented)", 32, AMD_RED, True)

caps = [
    ("RAG", "Local Knowledge Retrieval", "TF-IDF + BM25 + RRF\nPersistent disk index\nManual folder sync", GREEN),
    ("Tools", "Tool Invocation", "38+ tools\nAgentic loop (8 iter)\nBrowser, file, screen, keyboard", CYAN),
    ("Planning", "Multi-Step Task Planning", "AI decomposes tasks\nConditional tool calls\nSelf-correction on errors", GOLD),
    ("Memory", "Local Multi-Turn Memory", "JSONL persistence\nCross-session recall\nrecall / forget tools", RGBColor(0xBB, 0x33, 0xFF)),
    ("Privacy", "Permission & Privacy", "No hardcoded keys\nConfig file + env vars\nNo telemetry, local-first", AMD_RED),
]

x = Inches(0.3)
for tag, title, desc, color in caps:
    add_shape_with_text(slide, x, Inches(1.3), Inches(2.5), Inches(0.5),
                        tag, color, WHITE, 14)
    add_textbox(slide, x, Inches(2.0), Inches(2.5), Inches(0.5),
                title, 14, WHITE, True, PP_ALIGN.CENTER)
    add_textbox(slide, x, Inches(2.6), Inches(2.5), Inches(2.0),
                desc, 12, LIGHT_GRAY, False, PP_ALIGN.CENTER)
    x += Inches(2.6)

add_textbox(slide, Inches(0.5), Inches(5.5), Inches(12), Inches(0.5),
            "Minimum requirement: 2 of 5  →  Implemented: 5 of 5", 18, GREEN, True, PP_ALIGN.CENTER)

# ── Slide 5: Models & Deployment ──
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_textbox(slide, Inches(0.5), Inches(0.3), Inches(12), Inches(0.7),
            "Models & Local Deployment", 32, AMD_RED, True)

add_textbox(slide, Inches(0.5), Inches(1.2), Inches(6), Inches(0.5),
            "Models Used", 22, CYAN, True)

models = [
    ("DeepSeek-V4-Flash", "Primary LLM", "AMD Radeon Cloud API", "Fast inference, SSE streaming"),
    ("MiniCPM5-1B", "Fallback LLM", "AMD Radeon Cloud API", "Lightweight, auto-fallback"),
    ("Qwen3.6-35B-A3B", "Vision/Multimodal", "AMD Radeon Cloud API", "OCR, screen analysis"),
    ("Whisper Tiny EN", "Speech-to-Text", "100% Local (whisper.cpp)", "GGML quantized, on-device"),
]

y = Inches(1.8)
for name, role, deploy, opt in models:
    add_textbox(slide, Inches(0.5), y, Inches(3), Inches(0.4), name, 14, WHITE, True)
    add_textbox(slide, Inches(3.5), y, Inches(2.5), Inches(0.4), role, 14, GOLD)
    add_textbox(slide, Inches(6.0), y, Inches(3.5), Inches(0.4), deploy, 14, GREEN)
    add_textbox(slide, Inches(9.5), y, Inches(3.5), Inches(0.4), opt, 14, LIGHT_GRAY)
    y += Inches(0.6)

add_textbox(slide, Inches(0.5), Inches(4.5), Inches(12), Inches(0.5),
            "Full Local Deployment Path (AMD Radeon GPU + ROCm)", 22, CYAN, True)

steps = [
    "1. Deploy vLLM or llama.cpp (ROCm build) on AMD Radeon GPU",
    "2. Load model (e.g., DeepSeek-V4-Flash) and expose OpenAI-compatible API",
    "3. Point Argos to local endpoint via argos_config.txt:",
    "     server_url=http://localhost:8000/v1",
    "     api_key=local",
    "4. Whisper STT, RAG, Memory, and all tools already run 100% locally",
]
add_bullet_list(slide, Inches(0.5), Inches(5.2), Inches(12), Inches(2.0), steps, 15, LIGHT_GRAY)

# ── Slide 6: Optimization ──
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_textbox(slide, Inches(0.5), Inches(0.3), Inches(12), Inches(0.7),
            "Inference Speed Optimization on AMD Radeon GPU", 28, AMD_RED, True)

opts = [
    ("SSE Streaming", "Token-by-token streaming eliminates perceived latency.\nUsers see partial responses immediately."),
    ("Model Fallback Chain", "Primary (DeepSeek-V4-Flash) → Fallback (MiniCPM5-1B).\nPrevents total failure, always returns a response."),
    ("Agentic Loop Control", "Max 8 tool iterations, concise tool results,\nhistory limited to 20 messages to manage prompt size."),
    ("Local Component Optimization", "Whisper Tiny EN: ~75MB, sub-second on CPU.\nRAG persistent index: no re-indexing on startup.\nJSONL memory: O(1) append writes."),
    ("AMD Radeon GPU Path", "vLLM with ROCm: PagedAttention for KV cache.\nINT8/INT4 quantization via GGUF for llama.cpp.\nFlash Attention on ROCm for memory bandwidth."),
]

y = Inches(1.2)
for title, desc in opts:
    add_shape_with_text(slide, Inches(0.5), y, Inches(3.8), Inches(0.9),
                        title, RGBColor(0x0D, 0x47, 0x47), CYAN, 12)
    add_textbox(slide, Inches(4.5), y + Inches(0.05), Inches(8.3), Inches(0.9),
                desc, 14, LIGHT_GRAY)
    y += Inches(1.1)

# ── Slide 7: Tech Stack & Build ──
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_textbox(slide, Inches(0.5), Inches(0.3), Inches(12), Inches(0.7),
            "Tech Stack & Build", 32, AMD_RED, True)

add_textbox(slide, Inches(0.5), Inches(1.2), Inches(5.5), Inches(0.5),
            "Technology Stack", 22, CYAN, True)

stack = [
    "Language: C++17",
    "UI: Win32 API + Direct2D",
    "AI Inference: AMD Radeon Cloud API (OpenAI-compatible)",
    "Speech-to-Text: whisper.cpp (GGML, local)",
    "RAG Engine: Custom TF-IDF + BM25 + RRF",
    "Memory: JSONL persistent storage",
    "Browser Automation: Windows UI Automation API",
    "Build: MSBuild (Visual Studio 2022)",
    "Platform: Windows 10/11 x64",
]
add_bullet_list(slide, Inches(0.5), Inches(1.8), Inches(5.5), Inches(4.5), stack, 15, LIGHT_GRAY)

add_textbox(slide, Inches(6.5), Inches(1.2), Inches(6), Inches(0.5),
            "Build & Run", 22, CYAN, True)

build_steps = [
    "1. Install VS 2022 Build Tools (C++ workload)",
    "2. git clone https://github.com/Zrald1/amdagentic.git",
    "3. cd amdagentic\\cpp",
    "4. build.bat",
    "5. Output: build\\Debug\\argos.exe",
    "",
    "Configuration (3 methods):",
    "  - Environment variables: ARGOS_API_KEY",
    "  - Config file: argos_config.txt (next to exe)",
    "  - In-app Settings UI (auto-saves to config)",
    "",
    "No hardcoded API keys in source code.",
    "Zero external package manager dependencies.",
]
add_bullet_list(slide, Inches(6.5), Inches(1.8), Inches(6), Inches(5.0), build_steps, 14, LIGHT_GRAY)

# ── Slide 8: Demo Highlights / Thank You ──
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_textbox(slide, Inches(1), Inches(0.8), Inches(11), Inches(0.8),
            "Demo Highlights", 36, AMD_RED, True, PP_ALIGN.CENTER)

demos = [
    "Spartan warrior robot with shield & sword on desktop",
    "Electric lightning effect during tool execution",
    "Push-to-talk voice input (Caps Lock + Whisper)",
    "RAG folder sync & semantic search",
    "Browser automation: open URL → read screen → summarize",
    "Multi-turn conversation with persistent memory",
    "Streaming responses with reasoning/thoughts display",
]
add_bullet_list(slide, Inches(2), Inches(2.0), Inches(9), Inches(3.5), demos, 18, LIGHT_GRAY)

add_textbox(slide, Inches(1), Inches(5.8), Inches(11), Inches(0.5),
            "Thank You", 36, WHITE, True, PP_ALIGN.CENTER)
add_textbox(slide, Inches(1), Inches(6.5), Inches(11), Inches(0.5),
            "Team: GERALD BUSTILLA  |  Track 2  |  Argos — Agentic AI Companion",
            18, GOLD, False, PP_ALIGN.CENTER)

# Save
output_path = r"c:\Users\geral\amdagentic\docs\Argos_Presentation.pptx"
prs.save(output_path)
print(f"PPT saved to: {output_path}")
print(f"Total slides: {len(prs.slides)}")
